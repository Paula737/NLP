"""
Load + Extract stage of the ingestion pipeline.
Each function turns one kind of external knowledge source into a list of
LangChain Document objects (page_content + metadata).
"""
from langchain_core.documents import Document


def load_pdf(path: str):
    from pypdf import PdfReader
    reader = PdfReader(path)
    text = "\n".join((page.extract_text() or "") for page in reader.pages)
    return [Document(page_content=text, metadata={"source": path, "type": "pdf"})]


def load_docx(path: str):
    import docx
    d = docx.Document(path)
    text = "\n".join(p.text for p in d.paragraphs)
    return [Document(page_content=text, metadata={"source": path, "type": "docx"})]


def load_txt(path: str):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return [Document(page_content=text, metadata={"source": path, "type": "txt"})]


def load_code(path: str):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return [Document(page_content=text, metadata={"source": path, "type": "code"})]


def load_pptx(path: str):
    from pptx import Presentation
    prs = Presentation(path)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        parts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
        slides_text.append(f"Slide {i + 1}:\n" + "\n".join(parts))
    text = "\n\n".join(slides_text)
    return [Document(page_content=text, metadata={"source": path, "type": "pptx"})]


def load_webpage(url: str):
    import requests
    from bs4 import BeautifulSoup
    resp = requests.get(url, timeout=15, headers={"User-Agent": "Mozilla/5.0"})
    resp.raise_for_status()
    soup = BeautifulSoup(resp.text, "html.parser")
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()
    text = " ".join(soup.get_text(separator=" ").split())
    return [Document(page_content=text, metadata={"source": url, "type": "web"})]


def load_wikipedia(topic: str):
    import wikipedia
    page = wikipedia.page(topic, auto_suggest=False)
    return [Document(page_content=page.content, metadata={"source": page.url, "type": "wikipedia"})]


def load_audio(path: str):
    """Speech-to-text stage for WAV files, using faster-whisper."""
    from faster_whisper import WhisperModel
    model = WhisperModel("base", device="cpu", compute_type="int8")
    segments, _ = model.transcribe(path)
    text = " ".join(segment.text for segment in segments)
    return [Document(page_content=text, metadata={"source": path, "type": "audio"})]
